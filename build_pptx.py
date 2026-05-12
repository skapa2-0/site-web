"""Génère SKAPA_AI_Product_Platform_Pitch.pptx — présentation commerciale FR.
Identité SKAPA : indigo #4E56EC sur fond clair, accents orange #EA580C pour modules en dev.
Cible : équipes Produit / Tech (PM, CTO, lead dev).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ─── Brand tokens ────────────────────────────────────────────────────────────
PRIMARY      = RGBColor(0x4E, 0x56, 0xEC)
PRIMARY_DK   = RGBColor(0x3C, 0x44, 0xD4)
PRIMARY_LT   = RGBColor(0xEE, 0xF0, 0xFD)
INDIGO       = RGBColor(0x6C, 0x72, 0xE0)
INDIGO_LT    = RGBColor(0xA0, 0xA6, 0xF8)
ORANGE       = RGBColor(0xEA, 0x58, 0x0C)
ORANGE_LT    = RGBColor(0xFB, 0x92, 0x3C)
GREEN        = RGBColor(0x16, 0xA3, 0x4A)
BG           = RGBColor(0xF7, 0xF7, 0xFB)
BG2          = RGBColor(0xEE, 0xEE, 0xF8)
SURFACE      = RGBColor(0xFF, 0xFF, 0xFF)
TEXT         = RGBColor(0x0D, 0x0D, 0x1A)
TEXT_SEC     = RGBColor(0x60, 0x60, 0x80)
BORDER       = RGBColor(0xDC, 0xDD, 0xF5)
NAV_DARK     = RGBColor(0x13, 0x13, 0x1F)

FONT_DISPLAY = "Calibri"   # fallback proche de Syne en PPTX
FONT_BODY    = "Calibri Light"

# Slide 16:9 → 13.333 x 7.5 in
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ─── Helpers ─────────────────────────────────────────────────────────────────
def add_slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    bg_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg_shape.line.fill.background()
    bg_shape.fill.solid(); bg_shape.fill.fore_color.rgb = bg
    bg_shape.shadow.inherit = False
    return s

def add_rect(slide, x, y, w, h, fill=SURFACE, line=None, line_w=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        if line_w is not None: sh.line.width = line_w
    sh.shadow.inherit = False
    return sh

def add_round(slide, x, y, w, h, fill=SURFACE, line=None, radius=0.04):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = radius
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    sh.shadow.inherit = False
    return sh

def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT,
             font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             spacing=1.15, letter_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        if letter_spacing is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(int(letter_spacing * 100)))
    return tb

def add_eyebrow(slide, x, y, text, color=PRIMARY):
    """Petite étiquette tracking, type 'SECTION LABEL'."""
    return add_text(slide, x, y, Inches(6), Inches(0.3), text.upper(),
                    size=10, bold=True, color=color, letter_spacing=2.5)

def add_title(slide, x, y, w, h, text, *, size=40, color=TEXT, accent=None, accent_color=PRIMARY):
    """Titre Syne-like. Accent = morceau coloré (string contenue dans text)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.05
        if accent and accent in line:
            before, _, after = line.partition(accent)
            for txt, c in [(before, color), (accent, accent_color), (after, color)]:
                if not txt: continue
                r = p.add_run(); r.text = txt
                r.font.name = FONT_DISPLAY; r.font.size = Pt(size)
                r.font.bold = True; r.font.color.rgb = c
        else:
            r = p.add_run(); r.text = line
            r.font.name = FONT_DISPLAY; r.font.size = Pt(size)
            r.font.bold = True; r.font.color.rgb = color
    return tb

def add_pill(slide, x, y, text, *, fill=PRIMARY_LT, color=PRIMARY, w=Inches(1.4), h=Inches(0.32)):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.5
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT_BODY; r.font.size = Pt(9); r.font.bold = True
    r.font.color.rgb = color
    rPr = r._r.get_or_add_rPr(); rPr.set("spc", "200")
    return sh

def add_footer(slide, page_no, total):
    add_text(slide, Inches(0.6), Inches(7.15), Inches(6), Inches(0.25),
             "SKAPA · AI Product Platform", size=9, color=TEXT_SEC)
    add_text(slide, Inches(11.5), Inches(7.15), Inches(1.4), Inches(0.25),
             f"{page_no:02d} / {total:02d}", size=9, color=TEXT_SEC, align=PP_ALIGN.RIGHT)

def add_logo(slide, x=Inches(0.6), y=Inches(0.4)):
    add_text(slide, x, y, Inches(6), Inches(0.32),
             "SKAPA", size=14, bold=True, color=PRIMARY,
             font=FONT_DISPLAY, letter_spacing=3)
    add_text(slide, x, Inches(0.66), Inches(6), Inches(0.22),
             "AI PRODUCT PLATFORM", size=8, color=TEXT_SEC,
             letter_spacing=2.5)

TOTAL = 19   # nb total de slides

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(BG)
# Ambient indigo glow (top-right)
glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(-2), Inches(7), Inches(7))
glow.fill.solid(); glow.fill.fore_color.rgb = PRIMARY_LT
glow.line.fill.background(); glow.shadow.inherit = False
# Logo top-left
add_logo(s)
# Eyebrow
add_text(s, Inches(0.6), Inches(2.7), Inches(8), Inches(0.3),
         "AI PRODUCT PLATFORM · BY SKAPA", size=11, bold=True, color=PRIMARY, letter_spacing=3)
# Big title
add_title(s, Inches(0.6), Inches(3.1), Inches(12), Inches(2.6),
          "Le cycle produit entier.\nDans un seul système.",
          size=58, accent="Dans un seul système.", accent_color=PRIMARY)
# Sub
add_text(s, Inches(0.6), Inches(5.6), Inches(10), Inches(0.5),
         "De la conception au feedback utilisateur — connecté à votre release planning.",
         size=18, color=PRIMARY, bold=True)
add_text(s, Inches(0.6), Inches(6.05), Inches(10.5), Inches(0.7),
         "Quatre modules opérationnels. Un cycle continu. Une plateforme qui construit,\ncomprendre, audite et apprend — sans rupture entre les équipes.",
         size=12, color=TEXT_SEC, spacing=1.5)
add_text(s, Inches(11.5), Inches(7.15), Inches(1.4), Inches(0.25),
         "Mai 2026", size=9, color=TEXT_SEC, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — VISION
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(BG2)
add_logo(s)
add_eyebrow(s, Inches(0.6), Inches(1.5), "Vision")
add_title(s, Inches(0.6), Inches(1.85), Inches(12), Inches(2.5),
          "L'AI Product Platform est le premier système\nqui couvre l'intégralité du cycle produit —\nde la conception au feedback utilisateur.",
          size=30, accent="l'intégralité du cycle produit")
add_text(s, Inches(0.6), Inches(5.6), Inches(11.5), Inches(1),
         "Connectée au release planning Atlassian. Pensée pour les équipes Produit & Tech qui veulent\nfermer la boucle entre ce qu'elles construisent, ce que comprennent leurs utilisateurs, et ce\nqu'elles décident de prioriser ensuite.",
         size=14, color=TEXT_SEC, spacing=1.6)
add_footer(s, 2, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — LES 4 PILIERS
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(BG)
add_logo(s)
add_eyebrow(s, Inches(0.6), Inches(1.5), "Les 4 piliers")
add_title(s, Inches(0.6), Inches(1.85), Inches(12), Inches(1.4),
          "Construire. Comprendre. Auditer. Apprendre.", size=32)

pillars = [
    ("CONSTRUIRE", "Build OS",            "Full stack, dans vos contraintes,\ndéployé en production.",        PRIMARY,    "Beta"),
    ("COMPRENDRE", "Project OS",          "Communication avec le produit,\ncompréhension des choix.",          INDIGO_LT,  "Beta"),
    ("AUDITER",    "Health Care Product", "UX, accessibilité, éco-conception\nmesurées en continu.",            INDIGO,     "Beta"),
    ("APPRENDRE",  "Personae OS",         "Retours utilisateurs, personas\ndynamiques, release planning.",     ORANGE_LT,  "En dev"),
]
card_w = Inches(2.95); card_h = Inches(3.0); gap = Inches(0.1); x0 = Inches(0.6); y = Inches(3.7)
for i,(verb, name, desc, col, status) in enumerate(pillars):
    x = x0 + i*(card_w + gap)
    card = add_rect(s, x, y, card_w, card_h, fill=SURFACE, line=BORDER)
    # accent bar bottom
    add_rect(s, x, y+card_h-Emu(28000), card_w, Emu(28000), fill=col)
    add_text(s, x+Inches(0.3), y+Inches(0.3), card_w-Inches(0.6), Inches(0.3),
             verb, size=10, bold=True, color=col, letter_spacing=2.5)
    add_text(s, x+Inches(0.3), y+Inches(0.7), card_w-Inches(0.6), Inches(0.55),
             name, size=22, bold=True, color=TEXT, font=FONT_DISPLAY)
    add_text(s, x+Inches(0.3), y+Inches(1.45), card_w-Inches(0.6), Inches(1.2),
             desc, size=12, color=TEXT_SEC, spacing=1.4)
    pill_fill = ORANGE_LT if status=="En dev" else PRIMARY_LT
    pill_col  = ORANGE if status=="En dev" else PRIMARY
    add_pill(s, x+Inches(0.3), y+card_h-Inches(0.7), status, fill=pill_fill, color=pill_col, w=Inches(0.85))
add_footer(s, 3, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — LE CYCLE CONTINU
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(BG2)
add_logo(s)
add_eyebrow(s, Inches(0.6), Inches(1.5), "Le cycle continu")
add_title(s, Inches(0.6), Inches(1.85), Inches(12), Inches(1.4),
          "Un système vivant. Pas une suite d'outils.", size=30)

steps = [
    ("Build OS",         "Construit\n& déploie",   PRIMARY),
    ("Project OS",       "Comprend\n& documente",  INDIGO_LT),
    ("Health Care",      "Audite\n& recommande",   INDIGO),
    ("Personae OS",      "Apprend\n& priorise",    ORANGE_LT),
    ("Release Planning", "Atlassian",              GREEN),
]
y = Inches(4.3); step_w = Inches(2.0); gap_w = Inches(0.35); x0 = Inches(0.7)
for i,(name, sub, col) in enumerate(steps):
    x = x0 + i*(step_w + gap_w)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x+Inches(0.55), y, Inches(0.9), Inches(0.9))
    circle.fill.solid(); circle.fill.fore_color.rgb = col
    circle.line.fill.background(); circle.shadow.inherit = False
    add_text(s, x, y+Inches(1.05), step_w, Inches(0.4),
             name, size=14, bold=True, color=TEXT, font=FONT_DISPLAY, align=PP_ALIGN.CENTER)
    add_text(s, x, y+Inches(1.5), step_w, Inches(0.8),
             sub, size=11, color=TEXT_SEC, align=PP_ALIGN.CENTER, spacing=1.3)
    if i < len(steps)-1:
        arrow_x = x + step_w
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   arrow_x+Inches(0.04), y+Inches(0.4), gap_w-Inches(0.08), Inches(0.12))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = BORDER
        arrow.line.fill.background(); arrow.shadow.inherit = False

add_text(s, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.6),
         "Chaque module alimente le suivant. Les retours utilisateurs remontent dans le release planning.\nLe release planning alimente Build OS. Le cycle ne s'arrête jamais.",
         size=13, color=TEXT_SEC, align=PP_ALIGN.CENTER, spacing=1.5)
add_footer(s, 4, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — BUILD OS (focus)
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(BG)
add_logo(s)
add_pill(s, Inches(0.6), Inches(1.55), "MODULE 01 · BETA", fill=PRIMARY_LT, color=PRIMARY, w=Inches(1.7))
add_text(s, Inches(0.6), Inches(1.95), Inches(10), Inches(0.4),
         "CONSTRUIRE", size=11, bold=True, color=PRIMARY, letter_spacing=3)
add_title(s, Inches(0.6), Inches(2.3), Inches(12), Inches(1.4),
          "Build OS — full stack,\ndans vos contraintes.", size=40,
          accent="full stack,\ndans vos contraintes.")
add_text(s, Inches(0.6), Inches(4.6), Inches(11.5), Inches(0.8),
         "Génère votre produit digital simplement, sans paramétrage — dans l'environnement tech et\ndesign de chaque produit. Front, back, base de données. Déployé en production réelle.\nAccessible à n'importe quel métier — aussi simple que Lovable, infiniment plus puissant.",
         size=14, color=TEXT_SEC, spacing=1.55)
# Bullet features
bullets = [
    ("→ Génération full stack",       "dans votre DSM, vos conventions, votre architecture"),
    ("→ Accessible sans développeur", "aussi simple que Lovable"),
    ("→ Multi-produit",               "un contexte isolé par produit, un clic pour permuter"),
    ("→ Portail développeur",         "code auditable en 1 clic — Claude Code, Cursor, etc."),
    ("→ Intégration Jira",            "bidirectionnelle + documentation parallèle"),
    ("→ Déploiement réel",            "pas un prototype"),
]
col_x = [Inches(0.6), Inches(4.85), Inches(9.1)]
for i,(b, sub) in enumerate(bullets):
    cx = col_x[i % 3]; cy = Inches(6.2) if i < 3 else Inches(6.6)
    add_text(s, cx, cy, Inches(4.2), Inches(0.25), b, size=11, bold=True, color=PRIMARY)
    add_text(s, cx+Inches(0.18), cy+Inches(0.22), Inches(4.0), Inches(0.25), sub, size=10, color=TEXT_SEC)
add_footer(s, 5, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — BUILD OS · STATEMENT
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(NAV_DARK)
add_text(s, Inches(0.6), Inches(0.4), Inches(6), Inches(0.32),
         "SKAPA", size=14, bold=True, color=SURFACE,
         font=FONT_DISPLAY, letter_spacing=3)
add_text(s, Inches(0.6), Inches(0.66), Inches(6), Inches(0.22),
         "BUILD OS · STATEMENT", size=8, color=INDIGO_LT, letter_spacing=2.5)

add_title(s, Inches(0.8), Inches(2), Inches(12), Inches(2.3),
          "Les autres outils génèrent du code.\nBuild OS livre votre produit.",
          size=44, color=SURFACE, accent="Build OS livre votre produit.", accent_color=INDIGO_LT)
add_text(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2),
         "Dans votre design system, vos conventions, votre architecture.\nPas un prototype. Un produit.\n\nCo-construit avec des équipes métier qui connaissent déjà les outils du marché — et n'en\ntrouvaient aucun à la hauteur — et avec des développeurs lassés de retravailler du code généré.\nChaque feature répond à une friction réelle.",
         size=16, color=INDIGO_LT, spacing=1.6)
add_text(s, Inches(11.5), Inches(7.15), Inches(1.4), Inches(0.25),
         f"06 / {TOTAL:02d}", size=9, color=INDIGO_LT, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — BUILD OS · 4 AUDIENCES
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(BG)
add_logo(s)
add_eyebrow(s, Inches(0.6), Inches(1.5), "Pour chaque profil")
add_title(s, Inches(0.6), Inches(1.85), Inches(12), Inches(1.2),
          "Quatre audiences. Une seule plateforme.", size=30)

audiences = [
    ("01", "Métier & Product", PRIMARY,
     "Concevez comme si votre meilleure équipe de designers, devs et QA était à vos côtés.",
     ["Full stack livré en production — sans être développeur",
      "Dans votre DSM, vos contraintes et votre archi existante",
      "Multi-produit — un clic pour permuter",
      "Du concept au déploiement réel"]),
    ("02", "Développeurs", INDIGO,
     "Un système complet prêt à valider dans votre environnement. Pas une page blanche à réécrire.",
     ["Portail développeur — code contextualisé, documenté",
      "Compatible Claude Code, Cursor et votre infra",
      "Jira bidirectionnel — génération + lecture des tickets",
      "Passation fluide depuis le métier"]),
    ("03", "Architectes", INDIGO,
     "Dans vos contraintes. Pas à côté.",
     ["Architecture, sécurité, best practices respectées",
      "Patterns existants — pas de dette imposée",
      "Structure BDD existante intégrée nativement",
      "Workflow multi-agents — qualité AI Expert Squad"]),
    ("04", "Achats & Direction", INDIGO_LT,
     "Souveraineté. Conformité. Réduction des coûts.",
     ["LLM agnostic — vous gardez le contrôle",
      "On-premise ou cloud souverain EU",
      "Conforme aux exigences réglementaires",
      "Réduit le nombre d'outils dans votre stack"]),
]
card_w = Inches(6.05); card_h = Inches(1.95); gap = Inches(0.15); x0 = Inches(0.6); y0 = Inches(3.2)
for i,(num, role, col, title, items) in enumerate(audiences):
    cx = x0 + (i%2)*(card_w+gap); cy = y0 + (i//2)*(card_h+gap)
    card = add_rect(s, cx, cy, card_w, card_h, fill=SURFACE, line=BORDER)
    add_rect(s, cx, cy, Emu(38100), card_h, fill=col)  # left border
    add_text(s, cx+Inches(0.3), cy+Inches(0.15), Inches(0.6), Inches(0.3),
             num, size=11, bold=True, color=col, font=FONT_DISPLAY)
    add_text(s, cx+Inches(0.7), cy+Inches(0.17), Inches(3), Inches(0.3),
             role.upper(), size=10, bold=True, color=col, letter_spacing=2)
    add_text(s, cx+Inches(0.3), cy+Inches(0.5), card_w-Inches(0.5), Inches(0.55),
             title, size=12, bold=True, color=TEXT, font=FONT_DISPLAY, spacing=1.25)
    for j, it in enumerate(items):
        add_text(s, cx+Inches(0.3), cy+Inches(1.07)+j*Inches(0.21), card_w-Inches(0.5), Inches(0.22),
                 "·  " + it, size=9.5, color=TEXT_SEC)
add_footer(s, 7, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CE QUI REND BUILD OS UNIQUE (8 cartes)
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(BG)
add_logo(s)
add_eyebrow(s, Inches(0.6), Inches(1.5), "En résumé")
add_title(s, Inches(0.6), Inches(1.85), Inches(12), Inches(1.2),
          "Ce qui rend Build OS unique.", size=30)
recap = [
    ("01","Full stack réel",            "Front + back + BDD — un produit, pas un prototype."),
    ("02","Dans vos contraintes",       "DSM, conventions, archi BDD — votre code, pas du générique."),
    ("03","Sans développeur",           "Aussi simple que Lovable. Aussi puissant qu'une équipe senior."),
    ("04","Multi-produit",              "Un clic pour passer d'un produit à l'autre — contexte complet."),
    ("05","Connecté à votre SI",        "Jira, Figma, Git, APIs, BDD — votre stack, intacte."),
    ("06","Désilotage IT & métier",     "Continuité entre vision métier et exécution technique."),
    ("07","Code auditable en 1 clic",   "Le dev audite dans son environnement — zéro friction."),
    ("08","LLM agnostic & souverain",   "Cloud souverain ou on-premise selon vos règles."),
]
cw = Inches(2.95); ch = Inches(1.7); gap = Inches(0.15); x0 = Inches(0.6); y0 = Inches(3.3)
for i,(n,t,d) in enumerate(recap):
    cx = x0 + (i%4)*(cw+gap); cy = y0 + (i//4)*(ch+gap)
    featured = (n == "07")
    fill = PRIMARY_LT if featured else SURFACE
    add_rect(s, cx, cy, cw, ch, fill=fill, line=BORDER)
    add_text(s, cx+Inches(0.25), cy+Inches(0.2), Inches(0.6), Inches(0.3),
             n, size=14, bold=True, color=PRIMARY, font=FONT_DISPLAY)
    add_text(s, cx+Inches(0.25), cy+Inches(0.55), cw-Inches(0.4), Inches(0.4),
             t, size=13, bold=True, color=TEXT, font=FONT_DISPLAY)
    add_text(s, cx+Inches(0.25), cy+Inches(1.0), cw-Inches(0.4), Inches(0.65),
             d, size=10, color=TEXT_SEC, spacing=1.35)
add_footer(s, 8, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — MULTI-PRODUIT / MULTI-CONTEXTE
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(BG2)
add_logo(s)
add_eyebrow(s, Inches(0.6), Inches(1.5), "Multi-produit")
add_title(s, Inches(0.6), Inches(1.85), Inches(7), Inches(2),
          "Chaque produit,\nson propre univers.", size=36)
add_text(s, Inches(0.6), Inches(4.5), Inches(6.5), Inches(2),
         "Passez d'un produit à l'autre en un clic.\n\nBuild OS permute automatiquement le DSM,\nles conventions de code, la structure BDD\net les APIs connectées.\n\nChaque produit vit dans son propre contexte —\nisolé, cohérent, prêt.",
         size=14, color=TEXT_SEC, spacing=1.6)

# Visual: stacked product chips
vx = Inches(8.5); vy = Inches(2.4); vw = Inches(4.3)
add_round(s, vx, vy, vw, Inches(1.1), fill=PRIMARY_LT, line=PRIMARY, radius=0.15)
dot = s.shapes.add_shape(MSO_SHAPE.OVAL, vx+Inches(0.25), vy+Inches(0.42), Inches(0.18), Inches(0.18))
dot.fill.solid(); dot.fill.fore_color.rgb = GREEN; dot.line.fill.background(); dot.shadow.inherit = False
add_text(s, vx+Inches(0.55), vy+Inches(0.22), vw-Inches(0.7), Inches(0.32),
         "Un contexte par produit", size=13, bold=True, color=TEXT, font=FONT_DISPLAY)
add_text(s, vx+Inches(0.55), vy+Inches(0.55), vw-Inches(0.7), Inches(0.5),
         "DSM, code, BDD et APIs — chargés automatiquement", size=10, color=TEXT_SEC)

products = [
    ("● Produit A",  "React · Node.js · PostgreSQL · Figma"),
    ("● Produit B",  "Angular · Spring · Oracle · Jira"),
    ("● Produit C",  "Vue · FastAPI · MySQL · Figma"),
]
for i,(name, stack) in enumerate(products):
    cy = vy + Inches(1.4) + i*Inches(0.95)
    add_round(s, vx, cy, vw, Inches(0.82), fill=SURFACE, line=BORDER, radius=0.15)
    add_text(s, vx+Inches(0.3), cy+Inches(0.13), vw-Inches(0.4), Inches(0.3),
             name, size=12, bold=True, color=PRIMARY, font=FONT_DISPLAY)
    add_text(s, vx+Inches(0.3), cy+Inches(0.45), vw-Inches(0.4), Inches(0.3),
             stack, size=10, color=TEXT_SEC)
add_footer(s, 9, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — COMPARATIF
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(BG)
add_logo(s)
add_eyebrow(s, Inches(0.6), Inches(1.5), "Comparatif")
add_title(s, Inches(0.6), Inches(1.85), Inches(12), Inches(1.3),
          "Pensé pour libérer le métier et l'IT. Design. Ready to Host.",
          size=22, accent="Design. Ready to Host.")

headers = ["Critère", "Build OS", "Lovable", "Bolt", "Claude Code", "Claude Design"]
rows = [
    ("CÔTÉ MÉTIER", None),
    ("Génération full stack (front + back + BDD)",   ["✓","✗","◐","◐","✗"]),
    ("Dans votre DSM",                                ["✓","✗","✗","✗","◐"]),
    ("Contraintes front & structure BDD existante",  ["✓","✗","✗","✗","✗"]),
    ("Sans être développeur",                         ["✓","✓","✓","✗","✓"]),
    ("Multi-produit / multi-contexte",                ["✓","✗","✗","✗","✗"]),
    ("CÔTÉ DEV", None),
    ("Passation métier → dev",                        ["✓","✗","✗","✗","✗"]),
    ("Portail développeur + Jira bidirectionnel",    ["✓","✗","✗","◐","✗"]),
    ("Documentation parallèle + déploiement réel",   ["✓","✗","✗","✗","✗"]),
    ("INFRA IA", None),
    ("Workflow multi-agents",                         ["✓","✗","✗","✗","✗"]),
    ("LLM agnostic & souverain",                      ["✓","✗","✗","✗","✗"]),
]
tx = Inches(0.6); ty = Inches(3.4); col_w = [Inches(5.5), Inches(1.4), Inches(1.1), Inches(1.1), Inches(1.4), Inches(1.5)]
row_h = Inches(0.27)
# Header
cx = tx
for i,h in enumerate(headers):
    fill = PRIMARY if i==1 else BG2
    fc = SURFACE if i==1 else TEXT
    add_rect(s, cx, ty, col_w[i], Inches(0.32), fill=fill, line=BORDER)
    add_text(s, cx+Inches(0.1), ty+Inches(0.05), col_w[i]-Inches(0.1), Inches(0.25),
             h, size=10, bold=True, color=fc, align=PP_ALIGN.LEFT if i==0 else PP_ALIGN.CENTER)
    cx += col_w[i]
cy = ty + Inches(0.32)
for label, cells in rows:
    if cells is None:
        # section row
        add_rect(s, tx, cy, sum(col_w, Emu(0)), Inches(0.25), fill=PRIMARY_LT)
        add_text(s, tx+Inches(0.1), cy+Inches(0.03), Inches(10), Inches(0.2),
                 label, size=9, bold=True, color=PRIMARY, letter_spacing=2)
        cy += Inches(0.25); continue
    cx = tx
    add_rect(s, tx, cy, sum(col_w, Emu(0)), row_h, fill=SURFACE, line=BORDER)
    add_text(s, cx+Inches(0.1), cy+Inches(0.04), col_w[0]-Inches(0.1), row_h,
             label, size=10, color=TEXT)
    cx += col_w[0]
    for i,v in enumerate(cells):
        col_text = GREEN if v=="✓" else (ORANGE if v=="◐" else TEXT_SEC)
        bold = v=="✓"
        add_text(s, cx, cy+Inches(0.03), col_w[i+1], row_h,
                 v, size=12, bold=bold, color=col_text, align=PP_ALIGN.CENTER)
        cx += col_w[i+1]
    cy += row_h
add_footer(s, 10, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — INTÉGRATIONS / CONNECTEURS
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(BG2)
add_logo(s)
add_eyebrow(s, Inches(0.6), Inches(1.5), "Intégrations")
add_title(s, Inches(0.6), Inches(1.85), Inches(12), Inches(1.2),
          "Votre stack. Déjà là.", size=32)
add_text(s, Inches(0.6), Inches(2.95), Inches(11), Inches(0.5),
         "Build OS se connecte à vos outils existants dès le cadrage du pilote. Pas de migration, pas de friction — votre stack reste la vôtre.",
         size=12, color=TEXT_SEC)

cats = [
    ("Design",            ["Figma · Natif Dev Mode", "Storybook", "Design tokens"]),
    ("Gestion de projet", ["Jira Cloud/Server · Natif bidirectionnel", "Linear", "Notion"]),
    ("Versioning & CI/CD",["GitHub", "GitLab", "Azure DevOps", "Bitbucket", "Sonar · ESLint"]),
    ("Frameworks front",  ["React", "Angular", "Vue", "Next.js"]),
    ("Back-end & APIs",   ["Java · Spring", ".NET", "Python · FastAPI", "Node.js", "REST · GraphQL"]),
    ("Bases de données",  ["PostgreSQL", "SQL Server", "Oracle", "MongoDB", "MySQL · MariaDB"]),
    ("Déploiement",       ["Docker · K8s", "OpenShift", "AWS · Azure · GCP", "Cloud souverain EU", "On-premise"]),
    ("Sécurité & accès",  ["SSO · SAML · OIDC", "MFA", "Vault", "OWASP · SAST/DAST"]),
]
cw = Inches(3.0); ch = Inches(1.55); gap_x = Inches(0.15); gap_y = Inches(0.15); x0 = Inches(0.6); y0 = Inches(3.6)
for i,(cat, items) in enumerate(cats):
    cx = x0 + (i%4)*(cw+gap_x); cy = y0 + (i//4)*(ch+gap_y)
    add_rect(s, cx, cy, cw, ch, fill=SURFACE, line=BORDER)
    add_text(s, cx+Inches(0.2), cy+Inches(0.12), cw-Inches(0.3), Inches(0.28),
             cat.upper(), size=9, bold=True, color=PRIMARY, letter_spacing=2)
    add_text(s, cx+Inches(0.2), cy+Inches(0.45), cw-Inches(0.3), ch-Inches(0.5),
             "\n".join("·  "+it for it in items), size=9.5, color=TEXT, spacing=1.35)
add_footer(s, 11, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — ONBOARDING (4 étapes)
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(BG)
add_logo(s)
add_eyebrow(s, Inches(0.6), Inches(1.5), "Accompagnement")
add_title(s, Inches(0.6), Inches(1.85), Inches(12), Inches(1.5),
          "Un cadrage sur mesure\npour chaque organisation.", size=30)
add_text(s, Inches(0.6), Inches(3.7), Inches(5.5), Inches(2.5),
         "Chaque pilote commence par un atelier de cadrage —\npour que le code généré ressemble à votre production,\npas à un prototype générique.\n\nNous identifions ensemble votre stack, vos contraintes\net vos points bloquants avant de démarrer.\n\nUn document de synthèse SKAPA sert de base\ncontractuelle au kick-off.",
         size=12, color=TEXT_SEC, spacing=1.55)

steps2 = [
    ("01","Atelier de cadrage",                  "45 à 90 min · Métier + DSI · Stack, contraintes, succès."),
    ("02","Identification des points bloquants", "Réseau, Figma, données, RGPD — DSI, RSSI, DPO."),
    ("03","Prérequis validés & accord cadre",    "Document de synthèse SKAPA · base contractuelle."),
    ("04","Kick-off du pilote",                  "Build OS configuré sur votre contexte. Dès J1."),
]
bx = Inches(6.6); by = Inches(3.2); bw = Inches(6.2); bh = Inches(0.9); gap=Inches(0.08)
for i,(n,t,d) in enumerate(steps2):
    cy = by + i*(bh+gap)
    featured = (n=="04")
    fill = PRIMARY_LT if featured else SURFACE
    line = PRIMARY if featured else BORDER
    add_rect(s, bx, cy, bw, bh, fill=fill, line=line)
    nfill = PRIMARY if featured else BG2
    ncol  = SURFACE if featured else PRIMARY
    add_rect(s, bx+Inches(0.15), cy+Inches(0.15), Inches(0.6), Inches(0.6), fill=nfill, line=line)
    add_text(s, bx+Inches(0.15), cy+Inches(0.25), Inches(0.6), Inches(0.4),
             n, size=14, bold=True, color=ncol, font=FONT_DISPLAY, align=PP_ALIGN.CENTER)
    add_text(s, bx+Inches(0.9), cy+Inches(0.13), bw-Inches(1), Inches(0.32),
             t, size=13, bold=True, color=TEXT, font=FONT_DISPLAY)
    add_text(s, bx+Inches(0.9), cy+Inches(0.48), bw-Inches(1), Inches(0.4),
             d, size=10, color=TEXT_SEC)
add_footer(s, 12, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — PROJECT OS (module 02)
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(BG2)
add_logo(s)
add_pill(s, Inches(0.6), Inches(1.55), "MODULE 02 · BETA", fill=PRIMARY_LT, color=PRIMARY, w=Inches(1.7))
add_text(s, Inches(0.6), Inches(1.95), Inches(10), Inches(0.4),
         "COMPRENDRE", size=11, bold=True, color=INDIGO_LT, letter_spacing=3)
add_title(s, Inches(0.6), Inches(2.3), Inches(12), Inches(1.4),
          "Project OS — communiquez\navec votre produit.", size=38,
          accent="communiquez\navec votre produit.", accent_color=INDIGO)
add_text(s, Inches(0.6), Inches(4.5), Inches(11.5), Inches(0.9),
         "Interrogez votre produit comme un expert qui le connaît parfaitement — ses choix, son\nhistorique, ses fonctionnalités. Coordonnez vos équipes autour d'un contexte produit unique.",
         size=14, color=TEXT_SEC, spacing=1.55)
features = [
    "Interrogation produit en langage naturel",
    "Compréhension des choix de design et d'architecture",
    "Historique des décisions produit accessible",
    "Coordination des équipes autour du contexte produit",
]
for i,f in enumerate(features):
    add_text(s, Inches(0.6), Inches(5.7)+i*Inches(0.3), Inches(12), Inches(0.3),
             "→  " + f, size=13, color=TEXT)
add_footer(s, 13, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — HEALTH CARE PRODUCT (module 03)
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(BG)
add_logo(s)
add_pill(s, Inches(0.6), Inches(1.55), "MODULE 03 · BETA", fill=PRIMARY_LT, color=PRIMARY, w=Inches(1.7))
add_text(s, Inches(0.6), Inches(1.95), Inches(10), Inches(0.4),
         "AUDITER", size=11, bold=True, color=INDIGO, letter_spacing=3)
add_title(s, Inches(0.6), Inches(2.3), Inches(12), Inches(1.4),
          "Health Care Product —\nla qualité, mesurée en continu.", size=34,
          accent="la qualité, mesurée en continu.", accent_color=INDIGO)
add_text(s, Inches(0.6), Inches(4.5), Inches(11.5), Inches(0.9),
         "Mesurez et améliorez la qualité de votre produit en continu — UX, accessibilité et\néco-conception analysées automatiquement à chaque release.",
         size=14, color=TEXT_SEC, spacing=1.55)
audits = [
    ("UX",            "Audit UX automatisé à chaque release"),
    ("Accessibilité", "Conformité RGAA / WCAG AA"),
    ("Éco-conception","Empreinte numérique mesurée"),
    ("Reco IA",       "Recommandations actionnables générées"),
]
cw=Inches(2.95); ch=Inches(1.4); x0=Inches(0.6); y0=Inches(5.6); gap=Inches(0.15)
for i,(t,d) in enumerate(audits):
    cx = x0+i*(cw+gap)
    add_rect(s, cx, y0, cw, ch, fill=SURFACE, line=BORDER)
    add_text(s, cx+Inches(0.25), y0+Inches(0.2), cw-Inches(0.4), Inches(0.35),
             t, size=14, bold=True, color=INDIGO, font=FONT_DISPLAY)
    add_text(s, cx+Inches(0.25), y0+Inches(0.65), cw-Inches(0.4), Inches(0.7),
             d, size=10, color=TEXT_SEC, spacing=1.4)
add_footer(s, 14, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — PERSONAE OS + ATLASSIAN
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(BG2)
add_logo(s)
add_pill(s, Inches(0.6), Inches(1.55), "MODULE 04 · EN DÉVELOPPEMENT", fill=ORANGE_LT, color=ORANGE, w=Inches(2.6))
add_text(s, Inches(0.6), Inches(1.95), Inches(10), Inches(0.4),
         "APPRENDRE", size=11, bold=True, color=ORANGE, letter_spacing=3)
add_title(s, Inches(0.6), Inches(2.3), Inches(7.5), Inches(1.5),
          "Personae OS —\nfeedback → backlog → release.", size=28,
          accent="feedback → backlog → release.", accent_color=ORANGE)
add_text(s, Inches(0.6), Inches(4.7), Inches(7.5), Inches(2.2),
         "Transformez les retours utilisateurs en personas dynamiques\net en items de release planning — directement connecté à Atlassian.",
         size=13, color=TEXT_SEC, spacing=1.55)
feats = [
    "Collecte et analyse des retours utilisateurs",
    "Personas dynamiques mis à jour en continu",
    "Priorisation automatique des évolutions produit",
    "Intégration Atlassian — release planning alimenté",
    "Boucle feedback → backlog → release fermée",
]
for i,f in enumerate(feats):
    add_text(s, Inches(0.6), Inches(5.5)+i*Inches(0.27), Inches(7.5), Inches(0.27),
             "→  " + f, size=11, color=TEXT)

# Atlassian visual right
vx=Inches(8.6); vy=Inches(2.5); vw=Inches(4.2)
add_round(s, vx, vy, vw, Inches(4.4), fill=SURFACE, line=BORDER, radius=0.04)
add_text(s, vx+Inches(0.3), vy+Inches(0.2), vw-Inches(0.6), Inches(0.35),
         "📋  Release Planning — Atlassian", size=12, bold=True, color=TEXT, font=FONT_DISPLAY)
items = [
    ("Simplifier formulaire — 47 retours", "IA · P1",       PRIMARY,    PRIMARY_LT),
    ("Export PDF — demande fréquente",     "Users · P2",    GREEN,      RGBColor(0xE7,0xF7,0xEC)),
    ("Dashboard mobile — Agent terrain",   "Planifié · P3", INDIGO,     PRIMARY_LT),
    ("Intégration SSO — attente IT",       "Backlog",       TEXT_SEC,   BG2),
]
for i,(label, tag, dot, tag_bg) in enumerate(items):
    iy = vy + Inches(0.85) + i*Inches(0.78)
    add_round(s, vx+Inches(0.2), iy, vw-Inches(0.4), Inches(0.62), fill=BG2, line=BORDER, radius=0.15)
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, vx+Inches(0.35), iy+Inches(0.23), Inches(0.15), Inches(0.15))
    d.fill.solid(); d.fill.fore_color.rgb=dot; d.line.fill.background(); d.shadow.inherit=False
    add_text(s, vx+Inches(0.6), iy+Inches(0.12), vw-Inches(2.2), Inches(0.4),
             label, size=10, color=TEXT)
    add_pill(s, vx+vw-Inches(1.45), iy+Inches(0.15), tag, fill=tag_bg, color=dot, w=Inches(1.2), h=Inches(0.28))
add_footer(s, 15, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — ROADMAP
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(BG)
add_logo(s)
add_eyebrow(s, Inches(0.6), Inches(1.5), "Roadmap")
add_title(s, Inches(0.6), Inches(1.85), Inches(12), Inches(1.2),
          "Ce qui arrive.", size=32)
phases = [
    ("AUJOURD'HUI",  "Build OS · Project OS · Health Care",  "Les 3 premiers modules disponibles en Beta — opérationnels de manière indépendante.", "Beta",            GREEN,   PRIMARY),
    ("EN COURS",     "Personae OS",                          "Retours utilisateurs, personas dynamiques et intégration Atlassian release planning.","En développement", ORANGE, PRIMARY),
    ("PHASE 3",      "Connexion des modules",                "Les 4 modules interconnectés — partage de contexte, mémoire produit commune.",        "Roadmap",         PRIMARY, NAV_DARK),
    ("PHASE 4",      "Mémoire produit",                      "Contexte persistant et partagé entre tous les modules — historique, versions.",       "Roadmap",         PRIMARY, NAV_DARK),
]
y0=Inches(3.4); bh=Inches(0.85); gap=Inches(0.12)
for i,(phase, title, desc, status, status_col, phase_col) in enumerate(phases):
    cy = y0 + i*(bh+gap)
    # phase tag
    add_rect(s, Inches(0.6), cy, Inches(1.7), bh, fill=phase_col)
    add_text(s, Inches(0.6), cy+Inches(0.3), Inches(1.7), Inches(0.3),
             phase, size=10, bold=True, color=SURFACE, align=PP_ALIGN.CENTER, letter_spacing=2)
    # content
    add_rect(s, Inches(2.3), cy, Inches(10.4), bh, fill=SURFACE, line=BORDER)
    add_text(s, Inches(2.5), cy+Inches(0.13), Inches(3.2), Inches(0.3),
             title, size=12, bold=True, color=TEXT, font=FONT_DISPLAY)
    add_text(s, Inches(2.5), cy+Inches(0.45), Inches(7.5), Inches(0.4),
             desc, size=10, color=TEXT_SEC)
    bg_pill = RGBColor(0xE7,0xF7,0xEC) if status_col==GREEN else (
              RGBColor(0xFE,0xE6,0xD3) if status_col==ORANGE else PRIMARY_LT)
    add_pill(s, Inches(11.2), cy+Inches(0.28), status, fill=bg_pill, color=status_col, w=Inches(1.4), h=Inches(0.32))
add_footer(s, 16, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — 3 FORMATS D'ACCÈS
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(BG2)
add_logo(s)
add_eyebrow(s, Inches(0.6), Inches(1.5), "Formats d'accès")
add_title(s, Inches(0.6), Inches(1.85), Inches(12), Inches(1.4),
          "Trois façons de déployer la plateforme.", size=28)
formats = [
    ("FORMAT 01", "Meilleurs modèles",
     "La plateforme avec les meilleurs modèles du marché — veille continue, performances optimales.",
     ["Modèles les plus performants", "Mise à jour automatique", "Idéal pour démarrer rapidement"], False),
    ("FORMAT 02", "Modèles client",
     "Intégrez vos contrats LLM existants — Anthropic, Microsoft, Google et autres.",
     ["Vos contrats LLM valorisés", "Contrôle total sur le modèle", "Cohérence avec votre stack IA"], True),
    ("FORMAT 03", "Open Source & Souverain",
     "Modèles open source — cloud souverain EU ou on-premise sur votre infrastructure.",
     ["Zéro dépendance externe", "Données dans votre environnement", "Cloud souverain ou on-premise"], False),
]
cw=Inches(4.05); ch=Inches(3.6); x0=Inches(0.6); y0=Inches(3.4); gap=Inches(0.1)
for i,(tag, name, desc, items, featured) in enumerate(formats):
    cx = x0 + i*(cw+gap)
    fill = PRIMARY_LT if featured else SURFACE
    line = PRIMARY if featured else BORDER
    add_rect(s, cx, y0, cw, ch, fill=fill, line=line)
    add_text(s, cx+Inches(0.3), y0+Inches(0.3), cw-Inches(0.5), Inches(0.3),
             tag, size=10, bold=True, color=PRIMARY, letter_spacing=2.5)
    add_text(s, cx+Inches(0.3), y0+Inches(0.65), cw-Inches(0.5), Inches(0.5),
             name, size=18, bold=True, color=TEXT, font=FONT_DISPLAY)
    add_text(s, cx+Inches(0.3), y0+Inches(1.25), cw-Inches(0.5), Inches(1.0),
             desc, size=11, color=TEXT_SEC, spacing=1.5)
    for j,it in enumerate(items):
        add_text(s, cx+Inches(0.3), y0+Inches(2.4)+j*Inches(0.32), cw-Inches(0.5), Inches(0.3),
                 "→  " + it, size=11, color=TEXT)
add_footer(s, 17, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — PORTFOLIO (produits construits)
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(BG)
add_logo(s)
add_eyebrow(s, Inches(0.6), Inches(1.5), "Portfolio")
add_title(s, Inches(0.6), Inches(1.85), Inches(12), Inches(1.4),
          "Produits construits avec Build OS.", size=30)
add_text(s, Inches(0.6), Inches(3), Inches(11), Inches(0.45),
         "Des produits métiers réels, déployés en production — full stack, dans les contraintes des organisations.",
         size=12, color=TEXT_SEC)
portfolio = [
    ("Dashboard analytique",      "React · Node.js · PostgreSQL · Figma",     PRIMARY),
    ("Portail de gestion interne","Angular · Java Spring · Oracle · Jira",    INDIGO),
    ("Outil de pricing métier",   "Vue · Python FastAPI · MySQL · Figma",     INDIGO_LT),
    ("Votre prochain produit",    "Rejoignez la Beta —",                       ORANGE),
]
cw=Inches(3.0); ch=Inches(3.0); x0=Inches(0.6); y0=Inches(3.7); gap=Inches(0.13)
for i,(name, stack, col) in enumerate(portfolio):
    cx = x0 + i*(cw+gap)
    # preview
    add_round(s, cx, y0, cw, Inches(1.85), fill=NAV_DARK, line=None, radius=0.04)
    add_rect(s, cx+Inches(0.3), y0+Inches(0.3), cw-Inches(0.6), Inches(0.15), fill=col)
    add_rect(s, cx+Inches(0.3), y0+Inches(0.6), Inches(0.6), Inches(0.95), fill=RGBColor(0x2A,0x2D,0x52))
    add_rect(s, cx+Inches(1.0), y0+Inches(0.6), cw-Inches(1.3), Inches(0.25), fill=RGBColor(0x2A,0x2D,0x52))
    add_rect(s, cx+Inches(1.0), y0+Inches(0.95), cw-Inches(1.3), Inches(0.25), fill=RGBColor(0x2A,0x2D,0x52))
    add_rect(s, cx+Inches(1.0), y0+Inches(1.3),  cw-Inches(1.3), Inches(0.25), fill=RGBColor(0x2A,0x2D,0x52))
    # info
    add_rect(s, cx, y0+Inches(1.85), cw, ch-Inches(1.85), fill=SURFACE, line=BORDER)
    add_text(s, cx+Inches(0.2), y0+Inches(2.0), cw-Inches(0.4), Inches(0.4),
             name, size=13, bold=True, color=TEXT, font=FONT_DISPLAY)
    add_text(s, cx+Inches(0.2), y0+Inches(2.45), cw-Inches(0.4), Inches(0.5),
             stack, size=10, color=TEXT_SEC)
add_footer(s, 18, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19 — CTA / BETA
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide(PRIMARY)
# Subtle radial-like glow (white overlay oval)
glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3), Inches(0.5), Inches(8), Inches(7))
glow.fill.solid(); glow.fill.fore_color.rgb = PRIMARY_DK
glow.line.fill.background(); glow.shadow.inherit = False
# alpha not directly available, ok visually

add_text(s, Inches(0.6), Inches(0.5), Inches(8), Inches(0.4),
         "SKAPA · AI PRODUCT PLATFORM", size=12, bold=True, color=SURFACE, letter_spacing=3)

add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.4),
         "COMMENCER", size=11, bold=True, color=PRIMARY_LT, letter_spacing=3, align=PP_ALIGN.CENTER)
add_title(s, Inches(0.6), Inches(2.5), Inches(12), Inches(2.2),
          "Commencez par Build OS.", size=60, color=SURFACE,
          accent="Build OS.", accent_color=PRIMARY_LT)
add_text(s, Inches(2), Inches(4.6), Inches(9.3), Inches(0.7),
         "Le module phare de l'AI Product Platform.\nDisponible en Beta dès maintenant.",
         size=18, color=PRIMARY_LT, align=PP_ALIGN.CENTER, spacing=1.5)

# CTA buttons
btn1 = add_round(s, Inches(3.5), Inches(5.8), Inches(3), Inches(0.55), fill=SURFACE, radius=0.5)
add_text(s, Inches(3.5), Inches(5.95), Inches(3), Inches(0.3),
         "Rejoindre la Beta  →", size=13, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
btn2 = add_round(s, Inches(6.8), Inches(5.8), Inches(3), Inches(0.55), fill=PRIMARY_DK, line=PRIMARY_LT, radius=0.5)
add_text(s, Inches(6.8), Inches(5.95), Inches(3), Inches(0.3),
         "Démarrer un pilote", size=13, color=SURFACE, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
         "SKAPA · contact@skapa.io", size=10, color=PRIMARY_LT)
add_text(s, Inches(8), Inches(7.05), Inches(5), Inches(0.3),
         "© 2025 SKAPA · AI Product Platform", size=10, color=PRIMARY_LT, align=PP_ALIGN.RIGHT)

# ─── Save ────────────────────────────────────────────────────────────────────
out = "/home/user/site-web/SKAPA_AI_Product_Platform_Pitch.pptx"
prs.save(out)
print("✓ PPTX généré :", out)
print("  Slides :", len(prs.slides))
